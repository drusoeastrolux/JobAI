from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG    = RGBColor(0x0D, 0x0D, 0x0D)
ACCENT     = RGBColor(0x6C, 0x63, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xA0, 0xA0, 0xA0)
LIGHT_CARD = RGBColor(0x1A, 0x1A, 0x2E)

blank_layout = prs.slide_layouts[6]  # completely blank

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "WHAT I LEARNED", 1.0, 1.8, 11, 1.2, 52, bold=True, color=WHITE)
add_text(s, "Building with OpenRouter", 1.0, 3.1, 10, 0.8, 30, bold=False, color=ACCENT)
add_text(s, "An AI-powered career coach — model-agnostic by design", 1.0, 4.1, 10, 0.5, 16, color=GRAY, italic=True)
add_text(s, "Druso Escalante  ·  2026", 1.0, 6.5, 6, 0.4, 13, color=GRAY)

# ── SLIDE 2: The Problem ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "THE PROBLEM", 1.0, 0.5, 10, 0.7, 11, bold=True, color=ACCENT)
add_text(s, "Too Many\nAI Keys,\nNot Enough\nTime", 1.0, 1.2, 8, 3.5, 48, bold=True, color=WHITE)
bullets = [
    "→  Separate account for OpenAI, Anthropic, Google...",
    "→  Separate billing for each provider",
    "→  Separate SDK or API format per company",
    "→  Locked in the moment you pick one",
]
for i, b in enumerate(bullets):
    add_text(s, b, 1.0, 4.9 + i * 0.42, 11, 0.4, 14, color=GRAY)

# ── SLIDE 3: What is OpenRouter ───────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "THE SOLUTION", 1.0, 0.5, 10, 0.7, 11, bold=True, color=ACCENT)
add_text(s, "OpenRouter", 1.0, 1.1, 10, 1.0, 52, bold=True, color=WHITE)
add_text(s, "One API key. 100+ models. Zero lock-in.", 1.0, 2.2, 10, 0.5, 20, color=ACCENT)

cards = [
    ("ONE KEY", "Access GPT, Claude, Llama,\nGemma — all from one account"),
    ("FREE TIER", "Models ending in :free\ncost nothing to experiment with"),
    ("SAME FORMAT", "Uses the OpenAI SDK schema —\nno new library to learn"),
]
for i, (title, body) in enumerate(cards):
    x = 1.0 + i * 4.0
    add_rect(s, x, 3.2, 3.6, 2.8, LIGHT_CARD)
    add_text(s, title, x + 0.2, 3.4, 3.2, 0.5, 13, bold=True, color=ACCENT)
    add_text(s, body, x + 0.2, 4.0, 3.2, 1.5, 14, color=WHITE)

# ── SLIDE 4: How I Used It ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "HOW I USED IT", 1.0, 0.5, 10, 0.7, 11, bold=True, color=ACCENT)
add_text(s, "Three lines\nthat changed\neverything", 1.0, 1.1, 6, 2.5, 40, bold=True, color=WHITE)

add_rect(s, 7.2, 1.0, 5.6, 2.2, LIGHT_CARD)
code = 'const openai = new OpenAI({\n  baseURL: "https://openrouter.ai/api/v1",\n  apiKey: process.env.OPENROUTER_API_KEY\n})'
add_text(s, code, 7.4, 1.15, 5.2, 1.9, 11, color=ACCENT)

steps = [
    ("1", "Point the SDK at OpenRouter's URL instead of OpenAI"),
    ("2", "Set your model in a single .env variable"),
    ("3", "Change one line to swap the entire AI brain"),
]
for i, (num, txt) in enumerate(steps):
    y = 3.7 + i * 0.85
    add_rect(s, 1.0, y, 0.45, 0.45, ACCENT)
    add_text(s, num, 1.0, y, 0.45, 0.45, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, txt, 1.6, y + 0.04, 10, 0.4, 15, color=WHITE)

# ── SLIDE 5: Security Lesson ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "SECURITY", 1.0, 0.5, 10, 0.7, 11, bold=True, color=ACCENT)
add_text(s, "Never expose\nyour API key\nin the browser", 1.0, 1.1, 8, 2.8, 40, bold=True, color=WHITE)
add_text(s, "Anyone can open DevTools and steal a frontend API key.", 1.0, 4.0, 11, 0.5, 16, color=GRAY, italic=True)

cols = [
    ("WRONG", "Call OpenRouter\ndirectly from React\n\nKey visible in\nnetwork requests", RGBColor(0xFF,0x4C,0x4C)),
    ("RIGHT", "React  →  Your Express server\n→  OpenRouter\n\nKey never leaves\nyour machine", RGBColor(0x4C,0xFF,0x9A)),
]
for i, (label, body, col) in enumerate(cols):
    x = 1.0 + i * 5.8
    add_rect(s, x, 4.6, 5.2, 2.5, LIGHT_CARD)
    add_text(s, label, x + 0.2, 4.75, 4.8, 0.4, 13, bold=True, color=col)
    add_text(s, body, x + 0.2, 5.25, 4.8, 1.6, 14, color=WHITE)

# ── SLIDE 6: Streaming ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "STREAMING", 1.0, 0.5, 10, 0.7, 11, bold=True, color=ACCENT)
add_text(s, "The typing effect\nisn't magic —\nit's SSE", 1.0, 1.1, 7, 2.8, 40, bold=True, color=WHITE)
add_text(s, "Server-Sent Events (SSE)", 1.0, 4.1, 7, 0.5, 18, bold=True, color=ACCENT)
points = [
    "AI sends tokens one at a time as they're generated",
    "Server forwards each token to the browser instantly",
    "Browser appends each piece — looks like live typing",
    "This is exactly how ChatGPT works under the hood",
]
for i, p in enumerate(points):
    add_text(s, f"·  {p}", 1.0, 4.7 + i * 0.42, 11, 0.4, 14, color=GRAY)

# ── SLIDE 7: Key Takeaways ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "KEY TAKEAWAYS", 1.0, 0.5, 10, 0.7, 11, bold=True, color=ACCENT)
add_text(s, "What I'd tell\nsomeone starting\nfrom zero", 1.0, 1.1, 7, 2.5, 38, bold=True, color=WHITE)

takeaways = [
    ("Start free", "Use :free models to build. Upgrade when you're ready."),
    ("Stay agnostic", "Never hardcode a model name. Put it in .env."),
    ("Hide your keys", "Route through a server. Always."),
    ("Stream it", "Waiting for full responses feels broken. Stream instead."),
]
for i, (title, body) in enumerate(takeaways):
    y = 3.8 + i * 0.78
    add_text(s, f"{i+1}.  ", 1.0, y, 0.5, 0.5, 16, bold=True, color=ACCENT)
    add_text(s, title, 1.5, y, 2.5, 0.35, 15, bold=True, color=WHITE)
    add_text(s, body, 4.2, y + 0.02, 8.5, 0.35, 14, color=GRAY)

out = "/Users/drusoescalante/Desktop/AIBot/ChatAgent/OpenRouter_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
